"""Shared fixtures and helpers for the claw test suite.

All fixtures synthesize files into ``tmp_path`` so each test is hermetic.
Heavy fixtures (mp4, wav) call ``require_tool('ffmpeg')`` so missing tools fail
loudly instead of silently skipping — see _helpers.require_tool.
"""

from __future__ import annotations

import csv
import json
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Callable

import pytest
from click.testing import CliRunner

from ._helpers import require_tool


def pytest_configure(config: pytest.Config) -> None:
    config.addinivalue_line("markers", "network: requires network access.")
    config.addinivalue_line("markers", "slow: slower test, may exceed 1s.")
    config.addinivalue_line("markers", "flow: end-to-end multi-verb pipeline.")


def pytest_addoption(parser: pytest.Parser) -> None:
    parser.addoption("--no-network", action="store_true",
                     help="Skip tests marked with @pytest.mark.network.")


def pytest_collection_modifyitems(config: pytest.Config, items: list[pytest.Item]) -> None:
    if config.getoption("--no-network"):
        skip = pytest.mark.skip(reason="--no-network was passed")
        for item in items:
            if "network" in item.keywords:
                item.add_marker(skip)


# ──────────────────────────────────────────────────────────────────────────────
# CliRunner
# ──────────────────────────────────────────────────────────────────────────────

@pytest.fixture
def runner() -> CliRunner:
    return CliRunner()


# ──────────────────────────────────────────────────────────────────────────────
# Source-file fixtures (synthesized per test)
# ──────────────────────────────────────────────────────────────────────────────

@pytest.fixture
def sample_csv(tmp_path: Path) -> Callable[..., Path]:
    def _make(rows: int = 3, name: str = "data.csv",
              header: tuple[str, ...] = ("a", "b", "c")) -> Path:
        p = tmp_path / name
        with p.open("w", newline="", encoding="utf-8") as f:
            w = csv.writer(f)
            w.writerow(header)
            for i in range(rows):
                w.writerow([f"r{i}", i, i * 10])
        return p
    return _make


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "book.xlsx", sheet: str = "Data",
              rows: list[list] | None = None) -> Path:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet
        ws.append(["a", "b", "c"])
        for r in (rows if rows is not None else [[1, 2, 3], [4, 5, 6], [7, 8, 9]]):
            ws.append(r)
        p = tmp_path / name
        wb.save(p)
        return p
    return _make


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "doc.pdf", pages: int = 2,
              text: str = "hello world") -> Path:
        from reportlab.pdfgen import canvas
        p = tmp_path / name
        c = canvas.Canvas(str(p))
        for i in range(pages):
            c.drawString(72, 720, f"Page {i + 1} {text}")
            c.showPage()
        c.save()
        return p
    return _make


@pytest.fixture
def sample_pdf_multipage(tmp_path: Path) -> Callable[..., Path]:
    """5-page PDF with distinguishable per-page text — for split/merge tests."""
    def _make(name: str = "multi.pdf") -> Path:
        from reportlab.pdfgen import canvas
        p = tmp_path / name
        c = canvas.Canvas(str(p))
        for i in range(5):
            c.drawString(72, 720, f"PAGE_MARKER_{i + 1}")
            c.drawString(72, 700, "Lorem ipsum dolor sit amet, consectetur adipiscing elit.")
            c.showPage()
        c.save()
        return p
    return _make


@pytest.fixture
def sample_png(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "img.png", size: tuple[int, int] = (100, 100),
              color: tuple[int, int, int] = (200, 50, 50)) -> Path:
        from PIL import Image
        p = tmp_path / name
        im = Image.new("RGB", size, color)
        im.save(p, format="PNG")
        return p
    return _make


@pytest.fixture
def sample_jpg(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "img.jpg", size: tuple[int, int] = (200, 150),
              color: tuple[int, int, int] = (50, 150, 200)) -> Path:
        from PIL import Image
        p = tmp_path / name
        im = Image.new("RGB", size, color)
        im.save(p, format="JPEG", quality=85)
        return p
    return _make


@pytest.fixture
def sample_html(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "page.html",
              body: str = "<h1>T</h1><p>hi</p><a href='/x'>x</a>") -> Path:
        p = tmp_path / name
        p.write_text(f"<html><body>{body}</body></html>", encoding="utf-8")
        return p
    return _make


@pytest.fixture
def sample_html_rich(tmp_path: Path) -> Callable[..., Path]:
    """Realistic article HTML with headings, links, image, table — for web/html flows."""
    def _make(name: str = "article.html") -> Path:
        body = """
<html><head><title>The Article</title></head><body>
<header><nav><a href='/home'>Home</a></nav></header>
<article>
  <h1>Main Headline</h1>
  <p class="lead">A short lead paragraph that introduces the topic.</p>
  <h2>Section A</h2>
  <p>Body of section A with a <a href='https://example.com/abs'>link</a> and <a href='/rel/path'>relative</a>.</p>
  <table><thead><tr><th>k</th><th>v</th></tr></thead>
         <tbody><tr><td>x</td><td>1</td></tr><tr><td>y</td><td>2</td></tr></tbody></table>
  <img src='/img/photo.jpg' alt='photo'/>
  <h2>Section B</h2>
  <p>More body content. <script>alert('xss')</script><style>.bad{}</style></p>
</article>
<footer><p>(c) 2026</p></footer>
</body></html>"""
        p = tmp_path / name
        p.write_text(body, encoding="utf-8")
        return p
    return _make


@pytest.fixture
def sample_xml(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "doc.xml",
              body: str = "<root><a id='1'>v1</a><a id='2'>v2</a></root>") -> Path:
        p = tmp_path / name
        p.write_text(f'<?xml version="1.0"?>{body}', encoding="utf-8")
        return p
    return _make


@pytest.fixture
def sample_md(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "doc.md", body: str = "# T\n\nBody para.\n") -> Path:
        p = tmp_path / name
        p.write_text(body, encoding="utf-8")
        return p
    return _make


@pytest.fixture
def sample_md_rich(tmp_path: Path) -> Callable[..., Path]:
    """Markdown with headings, list, table, link, code-block — for from-md / convert flows."""
    def _make(name: str = "rich.md") -> Path:
        body = """# Top Heading

Intro paragraph with **bold** and a [link](https://example.com).

## Section A

- bullet one
- bullet two
- bullet three

## Section B

| k | v |
|---|---|
| x | 1 |
| y | 2 |

```python
print("hello")
```

> blockquote line.
"""
        p = tmp_path / name
        p.write_text(body, encoding="utf-8")
        return p
    return _make


@pytest.fixture
def sample_json_rows(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "rows.json",
              rows: list[dict] | None = None) -> Path:
        rows = rows or [{"a": 1, "b": "x"}, {"a": 2, "b": "y"}, {"a": 3, "b": "z"}]
        p = tmp_path / name
        p.write_text(json.dumps(rows), encoding="utf-8")
        return p
    return _make


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "deck.pptx", slides: int = 1) -> Path:
        import pptx
        prs = pptx.Presentation()
        for _ in range(slides):
            prs.slides.add_slide(prs.slide_layouts[0])
        p = tmp_path / name
        prs.save(p)
        return p
    return _make


@pytest.fixture
def sample_docx(tmp_path: Path) -> Callable[..., Path]:
    def _make(name: str = "doc.docx") -> Path:
        import docx
        d = docx.Document()
        d.add_heading("Title", level=1)
        d.add_paragraph("Body paragraph.")
        p = tmp_path / name
        d.save(p)
        return p
    return _make


@pytest.fixture
def sample_mp4(tmp_path: Path) -> Callable[..., Path]:
    """Synthesize a tiny mp4 via ffmpeg (REQUIRED — tests fail if missing)."""
    def _make(name: str = "clip.mp4", seconds: int = 2) -> Path:
        require_tool("ffmpeg")
        p = tmp_path / name
        subprocess.run(
            ["ffmpeg", "-y", "-loglevel", "error",
             "-f", "lavfi", "-i", f"color=c=blue:s=128x128:d={seconds}",
             "-f", "lavfi", "-i", f"sine=frequency=440:duration={seconds}",
             "-c:v", "libx264", "-pix_fmt", "yuv420p",
             "-c:a", "aac", "-shortest", str(p)],
            check=True,
        )
        return p
    return _make


@pytest.fixture
def sample_wav(tmp_path: Path) -> Callable[..., Path]:
    """1-second sine wave WAV via ffmpeg."""
    def _make(name: str = "tone.wav", seconds: int = 1, freq: int = 440) -> Path:
        require_tool("ffmpeg")
        p = tmp_path / name
        subprocess.run(
            ["ffmpeg", "-y", "-loglevel", "error",
             "-f", "lavfi", "-i", f"sine=frequency={freq}:duration={seconds}",
             str(p)],
            check=True,
        )
        return p
    return _make


@pytest.fixture
def sample_yaml_pipeline(tmp_path: Path) -> Callable[..., Path]:
    """Tiny YAML recipe — single shell step writing 'hello' — for pipeline tests."""
    def _make(name: str = "recipe.yaml", steps: list[dict] | None = None) -> Path:
        import yaml
        steps = steps or [
            {"name": "greet", "type": "shell", "cmd": f"{sys.executable} -c \"print('hello')\""},
        ]
        p = tmp_path / name
        p.write_text(yaml.safe_dump({"steps": steps}), encoding="utf-8")
        return p
    return _make
