"""Shared test helpers — invoke wrappers, assertions, artifact validators.

Use these everywhere instead of touching `runner.invoke(cli, [...])` directly.
Centralizing here means future test files stay short and consistent.

External-tool policy: ``require_tool(name)`` *fails* (does not skip) when the
tool is missing — we want CI / local runs to pressure the developer into
running ``healthcheck.py --install``. Tests that genuinely cannot run without
network (Drive / Gmail) still gate via ``--dry-run`` or ``--help`` only.
"""

from __future__ import annotations

import json
import shutil
import socket
import threading
from contextlib import contextmanager
from http.server import HTTPServer, SimpleHTTPRequestHandler
from pathlib import Path
from typing import Any, Iterator

from click.testing import CliRunner, Result

from claw.__main__ import cli


# ──────────────────────────────────────────────────────────────────────────────
# invocation
# ──────────────────────────────────────────────────────────────────────────────

def invoke(runner: CliRunner, *args: str | Path, input: str | None = None) -> Result:
    """Run ``claw <args...>`` and return the click ``Result``.

    Always asserts no Python traceback leaked through (ie. the CLI converted
    every error into a clean exit code + message). Caller still has to assert
    ``exit_code`` themselves.
    """
    argv = [str(a) for a in args]
    res = runner.invoke(cli, argv, input=input, catch_exceptions=True)
    assert "Traceback (most recent call last)" not in (res.output or ""), (
        f"unexpected traceback in `claw {' '.join(argv)}`:\n{res.output}"
    )
    return res


# ──────────────────────────────────────────────────────────────────────────────
# Subprocess-based invoke (for flow tests on Windows)
# ──────────────────────────────────────────────────────────────────────────────

class SubprocessResult:
    """Mimics enough of click's Result for our assert_* helpers."""

    def __init__(self, exit_code: int, output: str) -> None:
        self.exit_code = exit_code
        self.output = output


def invoke_subprocess(*args: str | Path, input: str | None = None,
                      retries: int = 5, retry_delay: float = 1.0) -> SubprocessResult:
    """Run ``claw <args>`` in a fresh subprocess.

    Slower than ``invoke()`` (~100ms cold-start per call) but releases all OS
    handles cleanly between calls — necessary for Windows flow tests that
    chain in-place writes on the same file. A small retry budget mitigates
    flaky Windows races (e.g. AV scanners briefly locking newly-written
    files between ``os.replace`` attempts).

    Use ``invoke()`` for unit tests; use ``invoke_subprocess()`` for flow
    tests that mutate the same artifact across stages.
    """
    import subprocess
    import sys
    import time

    argv = [sys.executable, "-m", "claw", *(str(a) for a in args)]
    last_output = ""
    last_code = 1
    for attempt in range(retries + 1):
        proc = subprocess.run(
            argv,
            input=input,
            capture_output=True,
            text=True,
            timeout=120,
            close_fds=True,
        )
        last_output = (proc.stdout or "") + (proc.stderr or "")
        last_code = proc.returncode
        # Retry only on transient Windows file-lock races.
        is_winlock = (
            last_code != 0
            and ("PermissionError" in last_output
                 or "WinError 5" in last_output
                 or "WinError 32" in last_output
                 or "Access is denied" in last_output)
        )
        if not is_winlock or attempt == retries:
            break
        time.sleep(retry_delay)
    if last_code == 0:
        assert "Traceback (most recent call last)" not in last_output, (
            f"unexpected traceback in `claw {' '.join(str(a) for a in args)}`:\n{last_output}"
        )
    return SubprocessResult(last_code, last_output)


# ──────────────────────────────────────────────────────────────────────────────
# assertions
# ──────────────────────────────────────────────────────────────────────────────

def assert_ok(res: Result, msg: str = "") -> None:
    """Assert exit_code == 0; surfaces full stdout on failure for debug-ability."""
    if res.exit_code != 0:
        raise AssertionError(
            f"{msg or 'expected exit 0'} — got exit {res.exit_code}\n"
            f"---- output ----\n{res.output}"
        )


def assert_exit(res: Result, code: int, msg: str = "") -> None:
    if res.exit_code != code:
        raise AssertionError(
            f"{msg or f'expected exit {code}'} — got exit {res.exit_code}\n"
            f"---- output ----\n{res.output}"
        )


def assert_json_output(res: Result) -> Any:
    """Assert ok + parseable JSON. Returns the parsed payload."""
    assert_ok(res)
    try:
        return json.loads(res.output)
    except json.JSONDecodeError as e:
        raise AssertionError(f"output not JSON: {e}\n---- output ----\n{res.output}")


# ──────────────────────────────────────────────────────────────────────────────
# external-tool gating  (FAIL, do not skip)
# ──────────────────────────────────────────────────────────────────────────────

def require_tool(name: str) -> None:
    """Hard-fail if ``name`` is not on PATH. Triggers `healthcheck.py --install`.

    We deliberately do NOT skip — silent skips let regressions ship.
    """
    if shutil.which(name) is None:
        raise AssertionError(
            f"external tool {name!r} missing from PATH. "
            f"Run: `python ~/.claude/skills/claude-claw/scripts/healthcheck.py --install`"
        )


# ──────────────────────────────────────────────────────────────────────────────
# artifact validators
# ──────────────────────────────────────────────────────────────────────────────

def assert_pdf_pages(path: Path, expected: int) -> None:
    import fitz  # PyMuPDF
    doc = fitz.open(path)
    try:
        if doc.page_count != expected:
            raise AssertionError(
                f"expected {expected} pages in {path.name}, got {doc.page_count}"
            )
    finally:
        doc.close()


def assert_pdf_contains(path: Path, needle: str) -> None:
    import fitz
    doc = fitz.open(path)
    try:
        text = "\n".join(p.get_text() for p in doc)
    finally:
        doc.close()
    if needle not in text:
        raise AssertionError(
            f"{needle!r} not in pdf text of {path.name}\n"
            f"---- first 200 chars ----\n{text[:200]}"
        )


def assert_xlsx_sheets(path: Path, names: list[str]) -> None:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True)
    try:
        actual = wb.sheetnames
        missing = [n for n in names if n not in actual]
        if missing:
            raise AssertionError(f"missing sheets {missing} in {path.name}; have {actual}")
    finally:
        # Critical on Windows: read_only mode keeps the underlying ZipFile mmap
        # alive until GC, which can prevent later os.replace() over the same path.
        wb.close()


def assert_xlsx_cell(path: Path, sheet: str, cell: str, expected: Any) -> None:
    import openpyxl
    wb = openpyxl.load_workbook(path)
    try:
        actual = wb[sheet][cell].value
        if actual != expected:
            raise AssertionError(
                f"{path.name}!{sheet}!{cell}: expected {expected!r}, got {actual!r}"
            )
    finally:
        wb.close()


def assert_docx_has_text(path: Path, needle: str) -> None:
    import docx
    d = docx.Document(path)
    text = "\n".join(p.text for p in d.paragraphs)
    if needle not in text:
        raise AssertionError(
            f"{needle!r} not in docx text of {path.name}\n"
            f"---- first 200 chars ----\n{text[:200]}"
        )


def assert_docx_heading(path: Path, text: str, level: int | None = None) -> None:
    import docx
    d = docx.Document(path)
    for p in d.paragraphs:
        if p.text == text and p.style.name.startswith("Heading"):
            if level is None:
                return
            if p.style.name == f"Heading {level}":
                return
    raise AssertionError(f"heading {text!r} (level={level}) not found in {path.name}")


def assert_pptx_slides(path: Path, expected: int) -> None:
    import pptx
    prs = pptx.Presentation(path)
    if len(prs.slides) != expected:
        raise AssertionError(
            f"expected {expected} slides in {path.name}, got {len(prs.slides)}"
        )


def assert_image_dims(path: Path, w: int, h: int) -> None:
    from PIL import Image
    with Image.open(path) as im:
        if (im.width, im.height) != (w, h):
            raise AssertionError(
                f"{path.name}: expected {w}x{h}, got {im.width}x{im.height}"
            )


def assert_image_format(path: Path, fmt: str) -> None:
    from PIL import Image
    with Image.open(path) as im:
        if (im.format or "").upper() != fmt.upper():
            raise AssertionError(
                f"{path.name}: expected format {fmt}, got {im.format}"
            )


def assert_video_duration(path: Path, expected: float, tol: float = 0.5) -> None:
    """Probe duration via ffprobe. Caller must require_tool('ffprobe') first."""
    import subprocess
    ff = shutil.which("ffprobe")
    assert ff, "ffprobe missing — call require_tool('ffprobe') first"
    out = subprocess.check_output(
        [ff, "-v", "error", "-show_entries", "format=duration",
         "-of", "default=nw=1:nk=1", str(path)],
        text=True,
    ).strip()
    actual = float(out)
    if abs(actual - expected) > tol:
        raise AssertionError(
            f"{path.name}: expected ~{expected}s (±{tol}), got {actual:.2f}s"
        )


# ──────────────────────────────────────────────────────────────────────────────
# local HTTP server (for `web` tests)
# ──────────────────────────────────────────────────────────────────────────────

@contextmanager
def local_http_server(directory: Path) -> Iterator[str]:
    """Spin up a SimpleHTTPRequestHandler bound to 127.0.0.1:<random>.

    Yields the base URL. Server tears down on exit. Single-threaded request
    handling is fine for the small fixture sets these tests use.
    """
    handler = lambda *a, **kw: SimpleHTTPRequestHandler(*a, directory=str(directory), **kw)  # noqa: E731
    sock = socket.socket()
    sock.bind(("127.0.0.1", 0))
    port = sock.getsockname()[1]
    sock.close()
    server = HTTPServer(("127.0.0.1", port), handler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    try:
        yield f"http://127.0.0.1:{port}"
    finally:
        server.shutdown()
        server.server_close()


# ──────────────────────────────────────────────────────────────────────────────
# YAML pipeline-recipe builder
# ──────────────────────────────────────────────────────────────────────────────

def make_yaml_pipeline(path: Path, steps: list[dict]) -> Path:
    """Write a minimal recipe to ``path``. ``steps`` is a list of dicts as the
    pipeline runner expects."""
    import yaml
    path.write_text(yaml.safe_dump({"steps": steps}), encoding="utf-8")
    return path
